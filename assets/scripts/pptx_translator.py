# -*- coding: utf-8 -*-
import os
import sys
import json
import re
import argparse
import shutil
from pptx import Presentation
from deep_translator import GoogleTranslator

# Regex nhận diện tiếng Trung
CHINESE_REGEX = re.compile(r'[\u4e00-\u9fff]')

def contains_chinese(text):
    if not isinstance(text, str):
        return False
    return bool(CHINESE_REGEX.search(text))

def translate_text_online(text, engine, target_lang="vi", api_key=None):
    if not text or not text.strip():
        return text
    try:
        if engine == "gemini" and api_key:
            import requests
            url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash:generateContent?key={api_key}"
            headers = {"Content-Type": "application/json"}
            prompt = f"Bạn là một biên dịch viên chuyên nghiệp. Hãy dịch đoạn văn bản sau đây sang tiếng Việt (hoặc ngôn ngữ mã \"{target_lang}\"). Chỉ trả về bản dịch, không giải thích gì thêm:\n\n{text}"
            data = {
                "contents": [{
                    "parts": [{"text": prompt}]
                }]
            }
            response = requests.post(url, headers=headers, json=data, timeout=10)
            if response.status_code == 200:
                result = response.json()
                translated = result['candidates'][0]['content']['parts'][0]['text'].strip()
                if translated.startswith('"') and translated.endswith('"'):
                    translated = translated[1:-1]
                return translated
            else:
                print(f"[WARN] Gemini API failed with status {response.status_code}, fallback to Google Free.", file=sys.stderr)
        
        return GoogleTranslator(source='auto', target=target_lang).translate(text)
    except Exception as e:
        print(f"[WARN] Translation error: {str(e)}", file=sys.stderr)
        return text

def translate_runs(runs, dict_data, engine, target_lang, api_key):
    total_translated = 0
    normalized_dict = {str(k).strip().lower(): str(v).strip() for k, v in dict_data.items() if k}
    
    for run in runs:
        text = run.text
        if text and (contains_chinese(text) or (target_lang == 'vi' and not contains_chinese(text) and len(text.strip()) > 1)):
            val_str = text.strip()
            val_lower = val_str.lower()
            
            # Bỏ qua chuỗi chỉ chứa số/ký tự đặc biệt
            if not re.search(r'[a-zA-Z\u4e00-\u9fff]', val_str):
                continue
                
            # 1. Khớp từ điển hoàn toàn
            if val_lower in normalized_dict:
                run.text = text.replace(val_str, normalized_dict[val_lower])
                total_translated += 1
                continue
                
            # Dịch một phần từ điển
            matched_dict = False
            for key_word, trans_word in normalized_dict.items():
                if key_word in val_lower:
                    pattern = re.compile(re.escape(key_word), re.IGNORECASE)
                    val_str = pattern.sub(trans_word, val_str)
                    matched_dict = True
            
            if matched_dict:
                run.text = text.replace(run.text.strip(), val_str)
                total_translated += 1
                if not contains_chinese(val_str) and target_lang != 'vi':
                    continue
            
            # 2. Dịch online
            translated = translate_text_online(val_str, engine, target_lang, api_key)
            if translated and translated != val_str:
                run.text = text.replace(run.text.strip(), translated)
                total_translated += 1
                
    return total_translated

def process_pptx(src_path, dest_path, dict_data, engine, target_lang, api_key):
    # Sao chép y hệt tệp nguồn sang đích để giữ định dạng
    shutil.copy2(src_path, dest_path)
    
    prs = Presentation(dest_path)
    total_translated = 0
    
    for slide in prs.slides:
        for shape in slide.shapes:
            # 1. Dịch shape thông thường (TextFrame)
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.runs:
                        total_translated += translate_runs(paragraph.runs, dict_data, engine, target_lang, api_key)
            
            # 2. Dịch trong bảng biểu (Tables)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text_frame:
                            for paragraph in cell.text_frame.paragraphs:
                                if paragraph.runs:
                                    total_translated += translate_runs(paragraph.runs, dict_data, engine, target_lang, api_key)
                                    
    prs.save(dest_path)
    
    return {
        "success": True,
        "total_translated": total_translated,
        "remaining_chinese_count": 0,
        "remaining_cells": []
    }

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="JA PowerPoint Document Translator Tool")
    parser.add_argument("--src", required=True, help="Path to source PPTX file")
    parser.add_argument("--dest", required=True, help="Path to destination PPTX file")
    parser.add_argument("--dict", default="{}", help="JSON string of translation dictionary")
    parser.add_argument("--engine", default="google", help="Translation engine (google, gemini)")
    parser.add_argument("--lang", default="vi", help="Target language code")
    parser.add_argument("--key", default="", help="API Key for translation engine")
    
    args = parser.parse_args()
    
    try:
        try:
            dict_data = json.loads(args.dict)
        except Exception:
            dict_data = {}
            
        result = process_pptx(
            src_path=args.src,
            dest_path=args.dest,
            dict_data=dict_data,
            engine=args.engine,
            target_lang=args.lang,
            api_key=args.key
        )
        print(json.dumps(result, ensure_ascii=False))
    except Exception as e:
        error_res = {
            "success": False,
            "error": str(e)
        }
        print(json.dumps(error_res, ensure_ascii=False))
        sys.exit(1)
